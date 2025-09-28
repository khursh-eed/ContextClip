# to get the job from db and convert tot ext using whisper ai + read slides

import asyncio
import time
import os
import json
import pandas as pd
import ffmpeg
from sqlalchemy.orm import Session
from typing import Optional, Dict, List
from pathlib import Path
import subprocess


from .database import SessionLocal, Job


async def process_job(job_id: str):
    db= SessionLocal()
    try:

        job = db.query(Job).filter(Job.id == job_id).first()
        if not job:
            print(f"Job {job_id} not found")
            return
        
        print(f"Starting processing for job {job_id}")
        job.status = "processing"
        db.commit()
        
        # Processing the audio file
        if job.media_path:
            try:
                # extract audio -> preprocess it -> transcribe + diarize -> give output
                print(f"Processing media file: {job.media_path}")
                processed_audio_path = await preprocess_audio(job.media_path, job_id)
                print(f"Audio preprocessing completed: {processed_audio_path}")
                
                transcript_data = await transcribe_and_diarize(processed_audio_path, job_id)
                print(f"Transcription completed, got {len(transcript_data.get('segments', []))} segments")
                
                # saving segments and transcript
                
                job_storage_path = f"storage/{job_id}"
                transcript_path = f"{job_storage_path}/transcript.json"
                segments_path = f"{job_storage_path}/segments.csv"
                
                with open(transcript_path, 'w', encoding='utf-8') as f:
                    json.dump(transcript_data, f, indent=2, ensure_ascii=False)
                print(f"Saved transcript to {transcript_path}")
                
                segments_df = create_segments_dataframe(transcript_data)
                segments_df.to_csv(segments_path, index=False)
                print(f"Saved segments to {segments_path}")
                
                # Process slides if available
                slides_dir = f"{job_storage_path}/slides"
                slide_links = {}
                
                if os.path.exists(slides_dir) and job.slides_count > 0:
                    print(f"Processing {job.slides_count} slides...")
                    
                    # Extract text from slides using OCR
                    extract_slides_from_file(job.slides_path ,slides_dir)
                    # Now slides_dir will contain images, ready for process_slides(slides_dir)
                    slide_texts = process_slides(slides_dir)
                    
                    if slide_texts:
                        # Link slides to transcript timestamps
                        slide_links = link_slides_to_transcript(
                            slide_texts, 
                            transcript_data.get('segments', [])
                        )
                        
                        # Save slide links
                        slide_links_path = f"{job_storage_path}/slide_links.json"
                        with open(slide_links_path, 'w', encoding='utf-8') as f:
                            json.dump(slide_links, f, indent=2, ensure_ascii=False)
                        print(f"Saved slide links to {slide_links_path}")
                        
                        print(f"Slide processing summary:")
                        for slide_id, link_info in slide_links.items():
                            timestamp = link_info.get('timestamp')
                            if timestamp is not None:
                                print(f"  Slide {slide_id}: {timestamp:.1f}s (confidence: {link_info.get('confidence_score', 0):.1f})")
                            else:
                                print(f"  Slide {slide_id}: No timestamp match found")
                
                # Update job with transcript path
                job.media_path = processed_audio_path  # Update to processed audio
                job.transcript_path = transcript_path
                db.commit()
                
                print(f"First 5 segments:")
                print(segments_df.head().to_string())
                
            except Exception as e:
                print(f"Error in media processing: {str(e)}")
                import traceback
                traceback.print_exc()
                raise
        
        # Mark as completed
        job.status = "done"
        db.commit()
        
        print(f"Job {job_id} completed successfully")
        
    except Exception as e:
        print(f"Error processing job {job_id}: {str(e)}")
        
        # Mark job as failed
        job = db.query(Job).filter(Job.id == job_id).first()
        if job:
            job.status = "error"
            db.commit()
    
    finally:
        db.close()


# now defining all the functions used

async def preprocess_audio(input_path: str, job_id: str) -> str:
    # need to convert audio to 16kHz - using ffmpeg

    try:
        job_storage_path = f"storage/{job_id}"
        output_path = f"{job_storage_path}/processed_audio.wav"
        os.makedirs(job_storage_path, exist_ok=True)
        
        
        print(f"Preprocessing audio: {input_path} -> {output_path}")
        
        # Use ffmpeg-python to convert audio
        stream = ffmpeg.input(input_path)
        stream = ffmpeg.output(
            stream,
            output_path,
            acodec='pcm_s16le',  # 16-bit PCM
            ac=1,                # Mono
            ar=16000            # 16kHz sample rate
        )
        
        # Run ffmpeg (overwrite output file)
        try:
            ffmpeg.run(stream, overwrite_output=True, quiet=True)
        except ffmpeg.Error as e:
            print('ffmpeg error:', e.stderr.decode())
            return input_path
        
        print(f"Audio preprocessing completed: {output_path}")
        return output_path
        
    except Exception as e:
        print(f"Error preprocessing audio: {str(e)}")
        # Return original path if preprocessing fails
        return input_path

async def transcribe_and_diarize(audio_path: str, job_id: str) -> Dict:
    # transcribing and diarzing 
    try:
        # trying WhisperX first -> if doesnt work using mock
        try:
            import whisperx
            print("Using WhisperX for transcription and diarization")
            return await transcribe_with_whisperx(audio_path, job_id)
        except Exception as e:
            print(f"WhisperX failed ({e}), not falling back to OpenAI Whisper for now")

            try:
                return await transcribe_with_openai_whisper(audio_path, job_id)
            except Exception as e:
                print(f"OpenAI Whisper failed ({e}), using mock transcription")
                return await transcribe_with_mock(audio_path, job_id)
            
    except Exception as e:
        print(f"Error in transcription: {str(e)}")
        return await transcribe_with_mock(audio_path, job_id)

async def transcribe_with_mock(audio_path: str, job_id: str) -> Dict:
    # generates fake segments based on audio duration
    try:
        import librosa
        print("Using mock transcription (for testing)")
        #audio duration
        # sr: sample rate (needs to be 16000)
        audio, sr = librosa.load(audio_path, sr=16000)
        duration = len(audio) / sr
        print(f"Audio duration: {duration:.2f} seconds")
        
        # random mock segments
        segments = []
        num_segments = max(1, int(duration / 2))  # One segment every 2 seconds
        mock_texts = [
            "This is a test transcription segment.",
            "The audio processing system is working correctly.",
            "Speaker diarization would separate different voices.",
            "This is just mock data for testing purposes.",
            "The real system would use WhisperX for actual transcription."
        ]
        
        for i in range(num_segments):
            start_time = i * 2.0
            end_time = min(start_time + 2.0, duration)
            speaker = f"SPEAKER_0{i % 2}"  # Alternate between two speakers
            text = mock_texts[i % len(mock_texts)]
            
            segments.append({
                "start": start_time,
                "end": end_time,
                "text": text,
                "speaker": speaker
            })
        
        return {
            "language": "en",
            "segments": segments,
            "job_id": job_id,
            "model": "mock-transcription",
            "audio_path": audio_path
        }
        
    except Exception as e:
        print(f"Mock transcription failed: {str(e)}")
        # Absolute fallback
        return {
            "language": "en",
            "segments": [{
                "start": 0.0,
                "end": 5.0,
                "text": "Mock transcription: Unable to process audio file.",
                "speaker": "SPEAKER_00"
            }],
            "job_id": job_id,
            "model": "fallback-mock",
            "audio_path": audio_path
        }
    
async def transcribe_with_whisperx(audio_path: str, job_id: str) -> Dict:
    try:
        import whisperx
        import torch

        # do not inlcude mps for mac, since whisperX doesnt support mps
        if torch.cuda.is_available():
            device = "cuda"
        else:
            device = "cpu"

        if device in ["cuda"]:
            compute_type = "float16"
        else:
            compute_type = "int8"
        
        # device = "cuda" if torch.cuda.is_available() else "cpu"
        # compute_type = "float16" if device == "cuda" else "int8"
        
        print(f"Loading WhisperX model on {device}")
        
        # 1. Transcribe with Whisper-small
        model = whisperx.load_model("small", device, compute_type=compute_type)
        audio = whisperx.load_audio(audio_path)
        # result = model.transcribe(audio, batch_size=16)
        result = model.transcribe(
            audio,
            batch_size=16,
            multilingual=True,
            max_new_tokens=128,
            clip_timestamps=None,
            hallucination_silence_threshold=0.1,
            hotwords=None
        )
        
        
        # 2. Align whisper output
        model_a, metadata = whisperx.load_align_model(language_code=result["language"], device=device)
        result = whisperx.align(result["segments"], model_a, metadata, audio, device, return_char_alignments=False)
        
        # 3. Assign speaker labels
        diarize_model = whisperx.DiarizationPipeline(use_auth_token=None, device=device)
        diarize_segments = diarize_model(audio_path)
        result = whisperx.assign_word_speakers(diarize_segments, result)
        
        return {
            "language": result.get("language", "unknown"),
            "segments": result["segments"],
            "job_id": job_id,
            "model": "whisperx-small",
            "audio_path": audio_path
        }
    except Exception as e:
        print(f"WhisperX failed: {str(e)}")
        raise
async def transcribe_with_openai_whisper(audio_path: str, job_id: str) -> Dict:
    """
    Use OpenAI Whisper for transcription (local model, no diarization)
    """
    try:
        import whisper
        
        print("Loading OpenAI Whisper model (small)")
        
        # Load the model (downloads on first use, then cached)
        model = whisper.load_model("small")
        
        print(f"Transcribing audio: {audio_path}")
        
        # Transcribe the audio
        result = model.transcribe(audio_path)
        
        # Convert to our format
        segments = []
        for segment in result.get("segments", []):
            segments.append({
                "start": segment.get("start", 0.0),
                "end": segment.get("end", 0.0),
                "text": segment.get("text", "").strip(),
                "speaker": "SPEAKER_00"  # No diarization in OpenAI Whisper
            })
        
        print(f"OpenAI Whisper completed: {len(segments)} segments")
        print(f"Detected language: {result.get('language', 'unknown')}")
        
        return {
            "language": result.get("language", "unknown"),
            "segments": segments,
            "job_id": job_id,
            "model": "openai-whisper-small",
            "audio_path": audio_path,
            "full_text": result.get("text", "")
        }
        
    except Exception as e:
        print(f"OpenAI Whisper failed: {str(e)}")
        raise


def create_segments_dataframe(transcript_data: Dict) -> pd.DataFrame:

    #convert segments into csv
    segments = []
    for segment in transcript_data.get("segments", []):
        segments.append({
            "start": segment.get("start", 0.0),
            "end": segment.get("end", 0.0),
            "speaker": segment.get("speaker", "SPEAKER_00"),
            "text": segment.get("text", "").strip()
        })
    
    return pd.DataFrame(segments)

def extract_slides_from_file(slide_file_or_dir: Optional[str], output_dir: str) -> None:
    """
    Prepares slide images (PNG) in output_dir from various input formats.
    Returns without error if input is empty or None.
    """
    import os
    if not slide_file_or_dir or not os.path.exists(slide_file_or_dir):
        print(f"No slides source found at '{slide_file_or_dir}', skipping slide extraction.")
        return
    if os.path.isdir(slide_file_or_dir):
        print(f"Slides already in image format in '{slide_file_or_dir}'.")
        return
    ext = os.path.splitext(slide_file_or_dir)[1].lower()
    if ext == ".pdf":
        convert_pdf_to_images(slide_file_or_dir, output_dir)
    elif ext in [".ppt", ".pptx"]:
        convert_ppt_to_images(slide_file_or_dir, output_dir)
    else:
        raise RuntimeError(f"Unsupported slide input format: {ext}")




def convert_pdf_to_images(pdf_path: str, slides_dir: str) -> list:
    """
    Converts each page of PDF to a PNG image and saves in slides_dir.
    Returns list of paths to images.
    """
    from pdf2image import convert_from_path
    import os

    os.makedirs(slides_dir, exist_ok=True)
    images = convert_from_path(pdf_path, dpi=200)
    img_paths = []

    for i, img in enumerate(images):
        img_path = os.path.join(slides_dir, f"slide_{i+1}.png")
        img.save(img_path, "PNG")
        img_paths.append(img_path)
    
    print(f"Converted PDF '{pdf_path}' to {len(img_paths)} slide images in '{slides_dir}'")
    return img_paths

def convert_ppt_to_images(ppt_path: str, slides_dir: str) -> list:
    """
    Try to convert PPT/PPTX to PNG images with LibreOffice headless.
    If that fails, fallback to text-only image rendering with python-pptx + Pillow.
    Returns list of image paths.
    """
    os.makedirs(slides_dir, exist_ok=True)
    
    # 1. Try LibreOffice conversion
    try:
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "png",
            "--outdir",
            slides_dir,
            ppt_path
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        # Collect generated images (usually slide0001.png, slide0002.png, etc)
        images = sorted([
            os.path.join(slides_dir, f) for f in os.listdir(slides_dir)
            if f.lower().endswith(".png")
        ])
        if images:
            print(f"LibreOffice converted PPT to {len(images)} slide images in '{slides_dir}'")
            return images
    
    except Exception as e:
        print(f"LibreOffice conversion failed: {e}")
    
    # 2. Fallback: use python-pptx + Pillow to create text-only images
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw
        
        prs = Presentation(ppt_path)
        img_paths = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = "\n".join([
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            ])
            img = Image.new('RGB', (1280, 720), color='white')
            draw = ImageDraw.Draw(img)
            draw.text((40, 40), slide_text[:5000], fill="black")
            img_path = os.path.join(slides_dir, f"slide_{i+1}.png")
            img.save(img_path, "PNG")
            img_paths.append(img_path)
        
        print(f"Fallback: Created {len(img_paths)} text-only slide images in '{slides_dir}'")
        return img_paths
    
    except Exception as e:
        print(f"Fallback text-only image generation failed: {e}")
        return []


def process_slides(slides_dir: str):

    # Process slide images using OCR to extract text
    # Returns: Dictionary mapping slide_id -> extracted text
    
    try:
        import pytesseract
        from PIL import Image
        import os
        
        # # Set Tesseract path for Windows
        # pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        
        if not os.path.exists(slides_dir):
            print(f"Slides directory not found: {slides_dir}")
            return {}
        
        slide_files = [f for f in os.listdir(slides_dir) 
                      if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff'))]
        
        if not slide_files:
            print("No slide images found")
            return {}
        
        print(f"Processing {len(slide_files)} slide images with OCR...")
        
        slide_texts = {}
        
        for i, slide_file in enumerate(sorted(slide_files)):
            slide_path = os.path.join(slides_dir, slide_file)
            text_path = os.path.join(slides_dir, f"slide_{i+1}.txt")
            slide_id = f"slide_{i+1}"
            
            try:
                # Open and process image with Tesseract
                image = Image.open(slide_path)
                
                # Extract text using OCR
                text = pytesseract.image_to_string(image, config='--psm 6')
                text = text.strip()
                
                # Save extracted text
                with open(text_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                
                slide_texts[slide_id] = text
                
                print(f"Slide {i+1}: Extracted {len(text)} characters")
                if text:
                    print(f"  Preview: {text[:100]}...")
                
            except Exception as e:
                print(f"Error processing slide {slide_file}: {str(e)}")
                # Create empty text file for failed slides
                with open(text_path, 'w', encoding='utf-8') as f:
                    f.write("")
                
                slide_texts[slide_id] = ""
        
        return slide_texts
        
    except ImportError:
        print("pytesseract not available, skipping slide processing")
        return {}
    except Exception as e:
        print(f"Error in slide processing: {str(e)}")
        return []

def link_slides_to_transcript(slide_texts: Dict[str, str], transcript_segments: List[Dict]) -> Dict:
    """
    Link slides to transcript timestamps using fuzzy matching
    """
    try:
        from rapidfuzz import fuzz, process
        import re
        
        if not slide_texts or not transcript_segments:
            return {}
        
        print(f"Linking {len(slide_texts)} slides to {len(transcript_segments)} transcript segments...")
        
        # Prepare transcript text for matching
        transcript_text_segments = []
        for segment in transcript_segments:
            # Clean text for better matching
            text = segment.get('text', '').strip()
            text = re.sub(r'\s+', ' ', text)  # Normalize whitespace
            text = re.sub(r'[^\w\s]', '', text)  # Remove punctuation
            
            if text:
                transcript_text_segments.append({
                    'text': text.lower(),
                    'timestamp': segment.get('start', 0.0),
                    'end': segment.get('end', 0.0),
                    'original_text': segment.get('text', '')
                })
        
        if not transcript_text_segments:
            print("No valid transcript segments for matching")
            return {}
        
        slide_links = {}
        
        for slide_id, slide_text in slide_texts.items():
            if not slide_text.strip():
                print(f"Slide {slide_id}: No text found, skipping")
                continue
            
            # Clean slide text for matching
            slide_text_clean = re.sub(r'\s+', ' ', slide_text)
            slide_text_clean = re.sub(r'[^\w\s]', '', slide_text_clean).lower()
            
            # Split slide text into phrases for better matching
            slide_phrases = [phrase.strip() for phrase in slide_text_clean.split('\n') 
                           if len(phrase.strip()) > 3]
            
            best_matches = []
            
            # Try to match each phrase
            for phrase in slide_phrases[:5]:  # Limit to first 5 phrases
                if len(phrase) < 5:  # Skip very short phrases
                    continue
                
                # Find best matching transcript segment
                matches = process.extract(
                    phrase,
                    [seg['text'] for seg in transcript_text_segments],
                    scorer=fuzz.partial_ratio,
                    limit=3
                )
                
                # Consider matches with score > 60
                for match in matches:
                    match_text, score, _ = match  # rapidfuzz returns (text, score, index)
                    if score > 60:
                        # Find the corresponding segment
                        for seg in transcript_text_segments:
                            if seg['text'] == match_text:
                                best_matches.append({
                                    'score': score,
                                    'timestamp': seg['timestamp'],
                                    'end': seg['end'],
                                    'phrase': phrase,
                                    'matched_text': seg['original_text']
                                })
                                break
            
            if best_matches:
                # Choose the match with highest score
                best_match = max(best_matches, key=lambda x: x['score'])
                
                slide_links[slide_id] = {
                    'timestamp': best_match['timestamp'],
                    'end_timestamp': best_match['end'],
                    'confidence_score': best_match['score'],
                    'matched_phrase': best_match['phrase'],
                    'matched_transcript': best_match['matched_text'],
                    'slide_text_preview': slide_text[:200]
                }
                
                print(f"Slide {slide_id}: Linked to {best_match['timestamp']:.1f}s "
                      f"(confidence: {best_match['score']:.1f})")
                print(f"  Phrase: '{best_match['phrase'][:50]}...'")
                print(f"  Matched: '{best_match['matched_text'][:50]}...'")
            else:
                print(f"Slide {slide_id}: No good matches found")
                # Store with no timestamp but keep the text
                slide_links[slide_id] = {
                    'timestamp': None,
                    'confidence_score': 0,
                    'slide_text_preview': slide_text[:200]
                }
        
        return slide_links
        
    except ImportError:
        print("rapidfuzz not available, skipping slide linking")
        return {}
    except Exception as e:
        print(f"Error in slide linking: {str(e)}")
        return {}

def generate_embeddings(text: str):
    """Stub for embedding generation"""
    # TODO: Implement sentence-transformers embeddings
    pass

def index_in_faiss(embeddings, metadata):
    """Stub for FAISS indexing"""
    # TODO: Implement FAISS vector indexing
    pass



    
